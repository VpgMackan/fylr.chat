"""Handler for pptx files."""

from pptx import Presentation
from io import BytesIO


def handle_pptx(file_content: bytes, **kwargs) -> str:
    """
    Extract text from pptx text files.

    Args:
        file_content: Raw file content as bytes
        **kwargs: Additional arguments (unused for text files)

    Returns:
        Extracted text as string
    """
    try:
        presentation = Presentation(BytesIO(file_content))
        full_text = []

        # Extract text from slides
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)

        return "\n".join(full_text)
    except Exception as e:
        raise ValueError(f"Failed to extract text from pptx: {e}")
